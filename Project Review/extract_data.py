"""
============================================================
  extract_data.py
  Academic Document Text Extractor
  ─────────────────────────────────
  Walks through the Files/ folder, extracts text from every
  PDF, PPTX and DOCX, infers course/module labels from the
  folder structure, and writes academic_documents_full.csv.

  Run once before project_da2.py.
  Usage:
      source venv/bin/activate
      python3 extract_data.py
============================================================
"""

import os
import re
import csv
import time
import traceback
from pathlib import Path

# ── Try importing extraction libraries ────────────────────────────────────────
try:
    import fitz          # PyMuPDF  — pip install PyMuPDF
    HAS_FITZ = True
except ImportError:
    HAS_FITZ = False
    print("⚠  PyMuPDF not found. PDF extraction disabled.")
    print("   Run: python3 -m pip install PyMuPDF")

try:
    from pptx import Presentation   # pip install python-pptx
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("⚠  python-pptx not found. PPTX extraction disabled.")

try:
    from docx import Document        # pip install python-docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    print("⚠  python-docx not found. DOCX extraction disabled.")

# ── Paths ─────────────────────────────────────────────────────────────────────
BASE_DIR    = Path(__file__).parent
FILES_DIR   = BASE_DIR / "Files"
OUTPUT_CSV  = BASE_DIR / "academic_documents_full.csv"

# ── Course code → short label mapping ─────────────────────────────────────────
COURSE_MAP = {
    "BMAT201L": "CVLA",
    "BCSE205L": "CAO",
    "BCSE308L": "CN",
    "BCSE302L": "DBMS",
    "BECE306L": "DCS",
    "BCSE202L": "DSA",
    "BECE102L": "DSD",
    "BECE301L": "DSP",
    "BCSE206L": "FDS",
    "BCSE401L": "IoT",
    "BCSE209L": "ML",
    "BECE204L": "MPMC",
    "BCSE303L": "OS",
    "BMAT202L": "PS",
    "BCSE304L": "TOC",
    "BCSE203E": "WP",
}

# ── Column names ──────────────────────────────────────────────────────────────
COLUMNS = [
    "filename", "relative_path",
    "course_code", "course_label",
    "module_label", "doc_type", "file_format",
    "word_count", "char_count",
    "extraction_status", "content"
]

# ── Helper: parse course code from folder name ─────────────────────────────────
def parse_course(folder_name: str):
    """Extract course code and label from folder name like 'DSA(BCSE202L)'."""
    m = re.search(r'\(([A-Z]+\d+[A-Z]?)\)', folder_name)
    if m:
        code = m.group(1)
        return code, COURSE_MAP.get(code, code)
    # fallback: use folder name as label
    label = folder_name.split("(")[0].strip().upper()
    return label, label

# ── Helper: parse module from path parts ──────────────────────────────────────
def parse_module(parts):
    """Search path parts for MODULE-N or MOD-N patterns."""
    patterns = [
        r'^MODULE[- _](\d+)$',
        r'^MOD[- _](\d+)$',
        r'^MODULE(\d+)$',
    ]
    for part in reversed(parts):   # search from deepest level up
        upper = part.upper().strip()
        for pat in patterns:
            m = re.match(pat, upper)
            if m:
                return f"module_{m.group(1)}"
    return "unknown"

# ── Helper: determine doc type from path ──────────────────────────────────────
DOC_TYPE_KEYWORDS = {
    "syllabus"       : "syllabus",
    "question paper" : "question_paper",
    "question papers": "question_paper",
    "qp"             : "question_paper",
    "cat-"           : "question_paper",
    "fat"            : "question_paper",
    "reference book" : "reference",
    "reference books": "reference",
    "book"           : "reference",
}

def parse_doc_type(parts, file_ext):
    """Infer document type from folder path and file extension."""
    path_str = " ".join(parts).lower()
    for kw, dtype in DOC_TYPE_KEYWORDS.items():
        if kw in path_str:
            return dtype
    # extension-based fallback
    if file_ext in (".pptx", ".ppt", ".pptm"):
        return "slides"
    if file_ext in (".pdf", ".docx", ".doc"):
        return "notes"
    return "other"

# ── Extraction functions ───────────────────────────────────────────────────────

def extract_pdf(path: Path) -> str:
    """Extract text from PDF using PyMuPDF."""
    if not HAS_FITZ:
        return ""
    try:
        doc = fitz.open(str(path))
        pages_text = []
        for page in doc:
            pages_text.append(page.get_text())
        doc.close()
        return " ".join(pages_text)
    except Exception:
        return ""

def extract_pptx(path: Path) -> str:
    """Extract text from PPTX using python-pptx."""
    if not HAS_PPTX:
        return ""
    try:
        prs = Presentation(str(path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
        return " ".join(texts)
    except Exception:
        return ""

def extract_docx(path: Path) -> str:
    """Extract text from DOCX using python-docx."""
    if not HAS_DOCX:
        return ""
    try:
        doc = Document(str(path))
        texts = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
        return " ".join(texts)
    except Exception:
        return ""

def extract_text(path: Path) -> tuple:
    """Dispatch to correct extractor. Returns (text, status)."""
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            text = extract_pdf(path)
        elif ext in (".pptx", ".pptm"):
            text = extract_pptx(path)
        elif ext in (".docx", ".doc"):
            text = extract_docx(path)
        else:
            return "", "skipped"

        text = clean_text(text)
        if len(text.split()) < 5:
            return "", "empty"
        return text, "success"
    except Exception:
        return "", "failed"

def clean_text(text: str) -> str:
    """Basic text cleaning."""
    # Remove excessive whitespace / newlines
    text = re.sub(r'\s+', ' ', text)
    # Remove non-printable characters
    text = re.sub(r'[^\x20-\x7E\u00C0-\u024F]', ' ', text)
    return text.strip()

# ── SKIP LIST: these folders contain non-academic content ─────────────────────
SKIP_FOLDERS = {
    "recordings", "youtube video links", "#cat-2 specific things",
    "#fat specific things"
}

SKIP_EXTENSIONS = {".jpg", ".jpeg", ".png", ".mp4", ".mp3", ".avi",
                   ".c", ".xlsx", ".xls", ".ppt", ".gif", ".bmp",
                   ".db", ".csv", ".txt", ".zip", ".7z", ".rar"}

# ── Main extraction loop ───────────────────────────────────────────────────────

def should_skip(path: Path) -> bool:
    """Return True if this file/folder should be skipped."""
    if path.suffix.lower() in SKIP_EXTENSIONS:
        return True
    # Check if any parent folder is in skip list
    for part in path.parts:
        if part.lower().strip() in SKIP_FOLDERS:
            return True
    if path.name.startswith("."):   # hidden files like .DS_Store
        return True
    return False

def process_files():
    if not FILES_DIR.exists():
        print(f"❌  Files/ folder not found at: {FILES_DIR}")
        print("    Make sure this script is in the same folder as the Files/ directory.")
        return

    all_files = [p for p in FILES_DIR.rglob("*") if p.is_file()]
    all_files = [p for p in all_files if not should_skip(p)]

    print(f"\n📂  Found {len(all_files)} processable files in Files/")
    print(f"📊  Extracting text — this will take a few minutes...\n")

    rows = []
    stats = {"success": 0, "empty": 0, "skipped": 0, "failed": 0}
    start = time.time()

    for i, filepath in enumerate(all_files):
        # ── Progress ────────────────────────────────────────────────────────────
        if i % 50 == 0 and i > 0:
            elapsed = time.time() - start
            eta = (elapsed / i) * (len(all_files) - i)
            print(f"  [{i}/{len(all_files)}]  success={stats['success']}  "
                  f"empty={stats['empty']}  failed={stats['failed']}  "
                  f"ETA {eta:.0f}s")

        # ── Parse path ──────────────────────────────────────────────────────────
        try:
            rel = filepath.relative_to(FILES_DIR)
            parts = list(rel.parts)

            if not parts:
                continue

            course_folder = parts[0]
            course_code, course_label = parse_course(course_folder)
            module_label  = parse_module(parts[1:])
            doc_type      = parse_doc_type(parts[1:], filepath.suffix.lower())
            file_format   = filepath.suffix.lower().lstrip(".")

            # ── Extract text ────────────────────────────────────────────────────
            text, status = extract_text(filepath)
            stats[status] = stats.get(status, 0) + 1

            wc = len(text.split()) if text else 0
            cc = len(text) if text else 0

            rows.append({
                "filename"        : filepath.name,
                "relative_path"   : str(rel),
                "course_code"     : course_code,
                "course_label"    : course_label,
                "module_label"    : module_label,
                "doc_type"        : doc_type,
                "file_format"     : file_format,
                "word_count"      : wc,
                "char_count"      : cc,
                "extraction_status": status,
                "content"         : text,
            })
        except Exception as e:
            stats["failed"] = stats.get("failed", 0) + 1
            continue

    # ── Write CSV ──────────────────────────────────────────────────────────────
    print(f"\n💾  Writing {len(rows)} records → {OUTPUT_CSV.name}")
    with open(OUTPUT_CSV, "w", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=COLUMNS)
        writer.writeheader()
        writer.writerows(rows)

    elapsed = time.time() - start
    successful = stats.get("success", 0)
    print(f"\n{'='*55}")
    print(f"  ✅  Extraction complete in {elapsed:.1f}s")
    print(f"  📄  Total files processed : {len(all_files)}")
    print(f"  ✔   Successfully extracted: {successful}")
    print(f"  ⚠   Empty / unreadable    : {stats.get('empty', 0)}")
    print(f"  ✗   Failed                : {stats.get('failed', 0)}")
    print(f"  ➜   Output: {OUTPUT_CSV}")
    print(f"{'='*55}")
    print(f"\n  ➜  Next: python3 project_da2.py")

if __name__ == "__main__":
    process_files()
